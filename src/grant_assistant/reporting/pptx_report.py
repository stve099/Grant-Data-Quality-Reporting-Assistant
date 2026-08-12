"""Generate a board-ready slide deck from the same calculated results.

Board meetings and funder check-ins run on slides, so today someone opens the
Word report and retypes the headline numbers into PowerPoint — the manual
retyping this project removes everywhere else.

The deck consumes the same :class:`ReportData` as the HTML and Word reports, so
a number can never differ between them: there is one source and four renderers
(HTML, PDF, Word, and this deck), and one profile controls the branding and the
section selection for all of them.
It is deliberately short. A deck is a summary; the report remains the artifact
with the detail, and slides that try to carry everything get read by nobody.

Requires the optional ``pptx`` extra. Charts additionally need the ``charts``
extra; without it the slides render with their tables and text intact.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Any

from grant_assistant.reporting.branding import brand_rgb, logo_bytes
from grant_assistant.reporting.chart_images import figure_png
from grant_assistant.reporting.context import ReportData
from grant_assistant.reporting.formatting import format_value as _fmt

logger = logging.getLogger(__name__)

# Tokens from docs/design_system.md, matching the Word report. BRAND is the
# fallback for the profile's brand_dark_color; the rest are not brandable.
BRAND = (0x1C, 0x5C, 0xAB)
INK = (0x1A, 0x1A, 0x1A)
MUTED = (0x52, 0x51, 0x4E)
GOOD = (0x1B, 0x7F, 0x4B)
BAD = (0xB4, 0x23, 0x1A)

#: 16:9, the only aspect ratio worth defaulting to for a screen or projector.
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


class PptxBackendError(Exception):
    """Raised when python-pptx is not installed."""


def _require_pptx() -> Any:
    try:
        import pptx
    except ImportError as exc:  # pragma: no cover - optional dependency
        raise PptxBackendError(
            "PowerPoint export requires the optional extra: uv sync --extra pptx"
        ) from exc
    return pptx


def write_pptx_report(report: ReportData, path: str | Path) -> Path:
    """Build the executive slide deck and write it to ``path``."""
    pptx = _require_pptx()
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = pptx.Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = presentation.slide_layouts[6]
    brand = brand_rgb(report)
    include = report.includes

    def add_slide() -> Any:
        return presentation.slides.add_slide(blank)

    def textbox(
        slide: Any,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        size: int = 18,
        bold: bool = False,
        color: tuple[int, int, int] = INK,
        align_center: bool = False,
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        if align_center:
            from pptx.enum.text import PP_ALIGN

            para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return frame

    def heading(slide: Any, text: str) -> None:
        textbox(slide, text, 0.6, 0.4, 12.1, 0.9, size=30, bold=True, color=brand)

    def bullets(slide: Any, items: list[str], top: float, size: int = 16) -> None:
        if not items:
            items = ["None noted."]
        frame = textbox(slide, items[0], 0.8, top, 11.7, 5.2, size=size)
        for item in items[1:]:
            para = frame.add_paragraph()
            run = para.add_run()
            run.text = item
            run.font.size = Pt(size)
            run.font.color.rgb = RGBColor(*INK)

    def chart_slide(title: str, key: str, charts: dict[str, Any]) -> bool:
        """A slide that is only worth creating if its chart renders."""
        figure = charts.get(key)
        if figure is None:
            return False
        png = figure_png(figure, width=1600, height=800)
        if png is None:
            return False
        slide = add_slide()
        heading(slide, title)
        slide.shapes.add_picture(io.BytesIO(png), Inches(1.0), Inches(1.4), width=Inches(11.3))
        return True

    a = report.analytics
    from grant_assistant.analytics.charts import standard_chart_set

    charts = standard_chart_set(a, report.audit)

    # --- Title --------------------------------------------------------------
    slide = add_slide()
    logo = logo_bytes(report)
    if logo is not None:
        # Centered above the title block; python-docx-style aspect scaling keeps it sane.
        slide.shapes.add_picture(
            io.BytesIO(logo[0]), Inches(SLIDE_WIDTH_IN / 2 - 1.0), Inches(1.1), width=Inches(2.0)
        )
    textbox(
        slide, report.title, 0.8, 2.4, 11.7, 1.2, size=44, bold=True, color=brand, align_center=True
    )
    textbox(
        slide,
        report.profile.grant_name,
        0.8,
        3.7,
        11.7,
        0.7,
        size=24,
        color=MUTED,
        align_center=True,
    )
    textbox(
        slide,
        f"{report.period_label}  ·  Prepared by {report.profile.report.prepared_by}",
        0.8,
        4.5,
        11.7,
        0.6,
        size=15,
        color=MUTED,
        align_center=True,
    )
    textbox(
        slide,
        "Synthetic demonstration data — contains no real client information.",
        0.8,
        5.2,
        11.7,
        0.5,
        size=12,
        color=MUTED,
        align_center=True,
    )

    # --- Headline numbers ---------------------------------------------------
    if include("population"):
        slide = add_slide()
        heading(slide, "At a glance")
        cards = [
            ("Households served", _fmt(a.households_served)),
            ("Individuals", _fmt(a.total_individuals)),
            ("Exits", _fmt(a.total_exits)),
            ("Permanent housing rate", _fmt(a.permanent_housing_rate, "percent")),
            ("Income increased", _fmt(a.pct_income_increased, "percent")),
            ("Follow-up completion", _fmt(a.overall_followup_completion_rate, "percent")),
        ]
        for index, (label, value) in enumerate(cards):
            column, row = index % 3, index // 3
            left = 0.8 + column * 4.0
            top = 1.6 + row * 2.3
            textbox(slide, value, left, top, 3.7, 1.0, size=40, bold=True, color=brand)
            textbox(slide, label, left, top + 1.0, 3.7, 0.6, size=14, color=MUTED)

    # --- Executive summary --------------------------------------------------
    if include("executive_summary"):
        slide = add_slide()
        heading(slide, "Executive summary")
        textbox(slide, report.executive_summary, 0.8, 1.5, 11.7, 5.0, size=15)

    # --- Charts -------------------------------------------------------------
    if include("programs"):
        chart_slide("Outcomes by program", "outcome_rates", charts)
    if include("measures"):
        chart_slide("Performance against target", "goal_vs_actual", charts)
    if include("enrollment"):
        chart_slide("Enrollment and exit trend", "enrollment_trends", charts)

    # --- Performance measures ----------------------------------------------
    if a.measures and include("measures"):
        slide = add_slide()
        heading(slide, "Performance measures")
        rows = len(a.measures) + 1
        table = slide.shapes.add_table(
            rows, 4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.4 * rows)
        ).table
        for index, label in enumerate(["Measure", "Target", "Actual", "Status"]):
            cell = table.cell(0, index)
            cell.text = label
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        for row_index, measure in enumerate(a.measures, start=1):
            met = measure.met
            values = [
                measure.name,
                _fmt(measure.target, measure.unit),
                _fmt(measure.actual, measure.unit),
                "Met" if met else ("Not met" if met is False else "n/a"),
            ]
            for column, value in enumerate(values):
                cell = table.cell(row_index, column)
                cell.text = value
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(12)
                if column == 3 and met is not None:
                    run.font.color.rgb = RGBColor(*(GOOD if met else BAD))
                    run.font.bold = True

    # --- Data quality -------------------------------------------------------
    if report.audit is not None and include("data_quality"):
        audit = report.audit
        slide = add_slide()
        heading(slide, "Data quality")
        textbox(
            slide, f"{audit.overall_score:.1f}", 0.8, 1.5, 3.0, 1.2, size=54, bold=True, color=brand
        )
        textbox(
            slide,
            f"Grade {audit.grade} · {audit.total_rows:,} records",
            0.8,
            2.8,
            3.6,
            0.6,
            size=14,
            color=MUTED,
        )
        blocking = [f"{i.rule_name} — {i.record_count} record(s)" for i in audit.blocking_issues]
        label = (
            "Blocking issues to resolve before submission:" if blocking else "No blocking issues."
        )
        textbox(slide, label, 4.8, 1.5, 7.7, 0.5, size=15, bold=True)
        bullets_items = blocking or ["This extract is ready to submit."]
        frame = textbox(slide, bullets_items[0], 4.8, 2.1, 7.7, 4.0, size=14)
        for item in bullets_items[1:]:
            para = frame.add_paragraph()
            run = para.add_run()
            run.text = item
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(*INK)

    if report.has_history and include("history"):
        history = report.history
        assert history is not None  # has_history
        slide = add_slide()
        heading(slide, "Data quality over time")
        movement = f"{history.since_previous:+.1f}" if history.since_previous is not None else "n/a"
        textbox(slide, movement, 0.8, 1.5, 3.0, 1.2, size=54, bold=True, color=brand)
        textbox(
            slide,
            f"points since the previous run · {history.runs} recorded",
            0.8,
            2.8,
            3.6,
            0.6,
            size=14,
            color=MUTED,
        )
        # The deck carries the same aging claim as the report, and the last few
        # runs rather than all of them: a slide holds a shape, not a ledger.
        lines = [f"{p.label}: {p.score:.1f}" for p in history.points[-5:]]
        lines += [f.describe() for f in history.persistent_findings[:3]]
        if history.resolved_rule_ids:
            lines.append("Resolved since last run: " + ", ".join(history.resolved_rule_ids[:6]))
        textbox(
            slide,
            "Recorded runs and long-standing findings",
            4.8,
            1.5,
            7.7,
            0.5,
            size=15,
            bold=True,
        )
        bullets(slide, lines[:8], 2.1, size=14)

    # --- Findings and actions ----------------------------------------------
    if include("findings"):
        slide = add_slide()
        heading(slide, "Key findings")
        bullets(slide, (report.insights.key_findings + report.insights.notable_trends)[:8], 1.5)

    if include("recommendations"):
        slide = add_slide()
        heading(slide, "Recommended actions")
        bullets(slide, report.insights.recommended_actions[:8], 1.5)

    # --- Methodology --------------------------------------------------------
    if include("methodology") or include("limitations"):
        slide = add_slide()
        heading(slide, "Methodology and limitations")
        mode = (
            "AI-assisted narrative grounded in deterministic calculations."
            if report.ai_generated_narrative
            else "Deterministic narrative (non-AI mode)."
        )
        items = []
        if include("methodology"):
            items += [
                "Every figure in this deck is calculated by the application; "
                "the AI layer only narrates.",
                mode,
            ]
        if include("limitations"):
            items += report.data_limitations()[:5]
        bullets(slide, items, 1.5, size=14)

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))
    logger.info("Wrote PowerPoint report to %s", path)
    return path
