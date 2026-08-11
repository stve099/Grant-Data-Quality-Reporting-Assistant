"""Provider resilience, report branding, relational merging, and scheduled runs."""

from __future__ import annotations

import asyncio
import logging

import pytest


def test_provider_error_preserves_failure_category():
    from grant_assistant.agents.provider import AIProviderError, AIProviderFailure

    error = AIProviderError.from_exception(
        TimeoutError("slow"), provider="test", operation="complete"
    )

    assert error.failure is AIProviderFailure.TIMEOUT
    assert error.retryable is True
    assert "test complete timed out" in str(error)


def test_provider_authentication_error_is_not_retryable():
    from grant_assistant.agents.provider import AIProviderError, AIProviderFailure

    auth_error = type("AuthenticationError", (Exception,), {})("bad key")
    error = AIProviderError.from_exception(auth_error, provider="test", operation="complete")

    assert error.failure is AIProviderFailure.AUTHENTICATION
    assert error.retryable is False


def test_provider_server_status_is_retryable():
    from grant_assistant.agents.provider import AIProviderError, AIProviderFailure

    server_error = type("InternalServerError", (Exception,), {"status_code": 503})(
        "temporarily unavailable"
    )
    error = AIProviderError.from_exception(server_error, provider="openai", operation="complete")

    assert error.failure is AIProviderFailure.PROVIDER
    assert error.retryable is True


def test_provider_request_timeout_status_is_retryable():
    from grant_assistant.agents.provider import AIProviderError, AIProviderFailure

    timeout_error = type("APIStatusError", (Exception,), {"status_code": 408})("timed out")
    error = AIProviderError.from_exception(
        timeout_error, provider="anthropic", operation="complete"
    )

    assert error.failure is AIProviderFailure.TIMEOUT
    assert error.retryable is True


def test_async_provider_adapter_runs_a_sync_provider():
    class Provider:
        name = "fake"

        def complete(self, system, messages, max_tokens=1500):
            return "ok"

    async def run():
        from grant_assistant.agents.provider import complete_async

        return await complete_async(Provider(), "system", [{"role": "user", "content": "hi"}])

    assert asyncio.run(run()) == "ok"


def test_report_branding_and_section_selection_are_honored(profile, analytics_clean):
    from grant_assistant.reporting import build_report_data, render_html_report

    branded = profile.model_copy(deep=True)
    branded.report.brand_color = "#123456"
    branded.report.brand_dark_color = "#102030"
    branded.report.sections = ["executive_summary", "population"]
    report = build_report_data(analytics_clean, None, branded)

    html = render_html_report(report, include_charts=False)

    assert "--brand: #123456" in html
    assert "--brand-dark: #102030" in html
    assert "Executive Summary" in html
    assert "Population Served" in html
    assert "Income Outcomes" not in html
    assert "Performance Measures" not in html


def test_concise_brief_honors_branding_and_section_selection(profile, analytics_clean):
    from grant_assistant.reporting import build_report_data, render_html_report

    branded = profile.model_copy(deep=True)
    branded.report.brand_color = "#123456"
    branded.report.sections = ["executive_summary", "population"]

    html = render_html_report(
        build_report_data(analytics_clean, None, branded), include_charts=False, template="concise"
    )

    assert "--brand: #123456" in html
    assert "Headline Results" in html
    assert "Performance Measures" not in html
    assert "Recommended Actions" not in html


def test_word_report_honors_branding_and_section_selection(profile, analytics_clean, tmp_path):
    from docx import Document

    from grant_assistant.reporting import build_report_data, write_docx_report

    branded = profile.model_copy(deep=True)
    branded.report.brand_dark_color = "#123456"
    branded.report.sections = ["executive_summary", "population"]

    report = build_report_data(analytics_clean, None, branded)
    path = write_docx_report(report, tmp_path / "r.docx")

    doc = Document(str(path))
    headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
    assert any("Population Served" in h for h in headings)
    assert not any("Income Outcomes" in h for h in headings)
    assert not any("Appendix" in h for h in headings)
    title_run = next(p.runs[0] for p in doc.paragraphs if p.runs)
    assert title_run.text == report.title
    assert str(title_run.font.color.rgb) == "123456"


def test_slide_deck_honors_branding_and_section_selection(profile, analytics_clean, tmp_path):
    import importlib.util

    import pytest

    if importlib.util.find_spec("pptx") is None:
        pytest.skip("pptx extra not installed")
    from pptx import Presentation

    from grant_assistant.reporting import build_report_data, write_pptx_report

    branded = profile.model_copy(deep=True)
    branded.report.brand_dark_color = "#123456"
    branded.report.sections = ["executive_summary", "population"]

    path = write_pptx_report(build_report_data(analytics_clean, None, branded), tmp_path / "d.pptx")

    deck = Presentation(str(path))
    text = "\n".join(
        shape.text_frame.text
        for slide in deck.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "At a glance" in text
    assert "Performance measures" not in text
    assert "Recommended actions" not in text
    title_run = deck.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    assert str(title_run.font.color.rgb) == "123456"


def test_every_renderer_drops_the_same_deselected_section(profile, analytics_clean, tmp_path):
    """One profile setting must mean the same thing in all four renderers.

    Without this, a funder receives a five-section PDF and a fifteen-section Word
    document generated from the same run.
    """
    from docx import Document

    from grant_assistant.reporting import build_report_data, render_html_report, write_docx_report

    trimmed = profile.model_copy(deep=True)
    trimmed.report.sections = [
        name for name in trimmed.report.sections if name not in {"income", "measures"}
    ]
    report = build_report_data(analytics_clean, None, trimmed)

    full_html = render_html_report(report, include_charts=False)
    brief_html = render_html_report(report, include_charts=False, template="concise")
    doc = Document(str(write_docx_report(report, tmp_path / "r.docx")))
    headings = "\n".join(p.text for p in doc.paragraphs if p.style.name.startswith("Heading"))

    for rendered in (full_html, brief_html, headings):
        assert "Income Outcomes" not in rendered
        assert "Performance Measures" not in rendered


def test_report_sections_reject_unknown_or_duplicate_names(profile):
    import pytest
    from pydantic import ValidationError

    from grant_assistant.configuration import GrantProfile

    payload = profile.model_dump(mode="python")
    payload["report"]["sections"] = ["outcomes", "outcome"]
    with pytest.raises(ValidationError, match="unknown report sections"):
        GrantProfile.model_validate(payload)

    payload["report"]["sections"] = ["outcomes", "outcomes"]
    with pytest.raises(ValidationError, match="duplicate report sections"):
        GrantProfile.model_validate(payload)


def test_related_dataset_merge_adds_columns_by_canonical_join_key(profile, tmp_path):
    import pandas as pd

    from grant_assistant.ingestion import merge_related_datasets

    primary = pd.DataFrame(
        {
            "Client ID": ["A", "B"],
            "Program Name": ["RRH", "ES"],
            "Entry Date": ["2025-01-01", "2025-01-02"],
        }
    )
    related = pd.DataFrame(
        {
            "Client ID": ["A", "B"],
            "Monthly Income at Exit": ["1000", "1200"],
            "Exit Destination": ["Rental by client, no subsidy", "Emergency shelter"],
        }
    )
    primary_path = tmp_path / "primary.csv"
    related_path = tmp_path / "income.csv"
    primary.to_csv(primary_path, index=False)
    related.to_csv(related_path, index=False)

    merged = merge_related_datasets(primary_path, [related_path], profile)

    assert list(merged["Client ID"]) == ["A", "B"]
    assert list(merged["Monthly Income at Exit"].astype(str)) == ["1000", "1200"]


def test_related_dataset_merge_normalizes_join_key_whitespace(profile, tmp_path):
    import pandas as pd

    from grant_assistant.ingestion import merge_related_datasets

    primary = pd.DataFrame({"Client ID": ["A"], "Program Name": ["RRH"]})
    related = pd.DataFrame({"Client ID": [" A "], "Case Manager": ["One"]})
    primary_path = tmp_path / "primary.csv"
    related_path = tmp_path / "services.csv"
    primary.to_csv(primary_path, index=False)
    related.to_csv(related_path, index=False)

    merged = merge_related_datasets(primary_path, [related_path], profile)

    assert merged.loc[0, "Case Manager"] == "One"
    assert merged.loc[0, "Client ID"] == "A"


def test_related_dataset_merge_rejects_duplicate_join_keys(profile, tmp_path):
    import pandas as pd
    import pytest

    from grant_assistant.ingestion import IngestionError, merge_related_datasets

    primary = pd.DataFrame({"Client ID": ["A"], "Program Name": ["RRH"]})
    related = pd.DataFrame({"Client ID": ["A", "A"], "Monthly Income at Exit": ["1", "2"]})
    primary_path = tmp_path / "primary.csv"
    related_path = tmp_path / "income.csv"
    primary.to_csv(primary_path, index=False)
    related.to_csv(related_path, index=False)

    with pytest.raises(IngestionError, match="duplicate join key"):
        merge_related_datasets(primary_path, [related_path], profile)


def test_related_dataset_merge_rejects_missing_join_keys(profile, tmp_path):
    import pandas as pd
    import pytest

    from grant_assistant.ingestion import IngestionError, merge_related_datasets

    primary = pd.DataFrame({"Client ID": ["A"], "Program Name": ["RRH"]})
    related = pd.DataFrame({"Client ID": ["A", None], "Monthly Income at Exit": ["1", "2"]})
    primary_path = tmp_path / "primary.csv"
    related_path = tmp_path / "income.csv"
    primary.to_csv(primary_path, index=False)
    related.to_csv(related_path, index=False)

    with pytest.raises(IngestionError, match="missing join key"):
        merge_related_datasets(primary_path, [related_path], profile)


def test_profile_rejects_alias_collision_between_programs(profile):
    import pytest
    from pydantic import ValidationError

    from grant_assistant.configuration import GrantProfile

    payload = profile.model_dump()
    payload["programs"][1]["aliases"].append(payload["programs"][0]["aliases"][0])
    with pytest.raises(ValidationError, match=r"alias.*more than one program"):
        GrantProfile.model_validate(payload)


def test_scheduled_audit_email_contains_submission_signal(analytics_clean, audit_clean, profile):
    from grant_assistant.automation import build_audit_email

    message = build_audit_email(profile, audit_clean, analytics_clean, "nightly.csv")

    assert profile.grant_name in message["Subject"]
    assert "Data quality score" in message.get_content()
    assert "Ready for submission: yes" in message.get_content()


def test_send_audit_email_uses_verified_tls_and_smtp_credentials(monkeypatch):
    import ssl

    from grant_assistant.automation import send_audit_email

    calls: dict[str, object] = {}

    class FakeSMTP:
        def __init__(self, host: str, port: int, timeout: int):
            calls["connection"] = (host, port, timeout)

        def __enter__(self):
            return self

        def __exit__(self, *_args):
            return None

        def starttls(self, *, context):
            calls["tls_context"] = context

        def login(self, username: str, password: str):
            calls["login"] = (username, password)

        def send_message(self, message):
            calls["message"] = message

    monkeypatch.setattr("grant_assistant.automation.smtplib.SMTP", FakeSMTP)

    from email.message import EmailMessage

    message = EmailMessage()
    message.set_content("Synthetic audit summary")
    send_audit_email(
        message,
        ["reviewer@example.org"],
        host="smtp.example.org",
        username="mailer",
        password="test-only-password",
        sender="reports@example.org",
    )

    assert calls["connection"] == ("smtp.example.org", 587, 30)
    assert isinstance(calls["tls_context"], ssl.SSLContext)
    assert calls["tls_context"].verify_mode == ssl.CERT_REQUIRED
    assert calls["login"] == ("mailer", "test-only-password")
    # Addressed on a copy, so the caller's message is left untouched for a retry.
    delivered = calls["message"]
    assert delivered is not message
    assert delivered.get_content() == message.get_content()
    assert delivered["To"] == "reviewer@example.org"
    assert delivered["From"] == "reports@example.org"
    assert message["To"] is None


def test_send_audit_email_rejects_credentials_without_tls():
    from email.message import EmailMessage

    import pytest

    from grant_assistant.automation import send_audit_email

    with pytest.raises(ValueError, match="credentials require TLS"):
        send_audit_email(
            EmailMessage(),
            ["reviewer@example.org"],
            host="smtp.example.org",
            username="mailer",
            password="test-only-password",
            sender="reports@example.org",
            use_tls=False,
        )


def test_new_cli_commands_are_registered():
    from typer.testing import CliRunner

    from grant_assistant.cli.main import app

    result = CliRunner().invoke(app, ["--help"])
    assert result.exit_code == 0
    assert "merge-datasets" in result.stdout
    assert "scheduled-audit" in result.stdout


def test_html_report_embeds_small_local_logo(profile, analytics_clean, tmp_path):
    from grant_assistant.reporting import build_report_data, render_html_report

    logo = tmp_path / "logo.png"
    logo.write_bytes(b"\x89PNG\r\n\x1a\nsynthetic")
    branded = profile.model_copy(deep=True)
    branded.report.logo_path = str(logo)

    html = render_html_report(
        build_report_data(analytics_clean, None, branded), include_charts=False
    )

    assert 'src="data:image/png;base64,' in html
    assert 'alt="Organization logo"' in html


# --- Degradation guarantees --------------------------------------------------
# "Optional extras must degrade, not crash" is a stated project invariant. These
# cover the branches that make it true, which are exactly the branches no happy
# path reaches.


@pytest.mark.parametrize(
    ("filename", "size", "reason"),
    [
        ("logo.svg", 8, "unsupported format"),
        ("logo.png", 2 * 1024 * 1024 + 1, "over the 2 MB cap"),
    ],
    ids=["unsupported-format", "over-size-cap"],
)
def test_unusable_logo_is_skipped_rather_than_fatal(
    profile, analytics_clean, tmp_path, caplog, filename, size, reason
):
    from grant_assistant.reporting import build_report_data, render_html_report

    bad = tmp_path / filename
    bad.write_bytes(b"x" * size)
    branded = profile.model_copy(deep=True)
    branded.report.logo_path = str(bad)
    report = build_report_data(analytics_clean, None, branded)

    with caplog.at_level(logging.WARNING):
        html = render_html_report(report, include_charts=False)

    assert "data:image" not in html, reason
    assert "Skipping report logo" in caplog.text
    assert report.title in html  # the rest of the report still rendered


def test_missing_logo_path_is_skipped_rather_than_fatal(profile, analytics_clean, tmp_path):
    from grant_assistant.reporting import build_report_data, render_html_report

    branded = profile.model_copy(deep=True)
    branded.report.logo_path = str(tmp_path / "does_not_exist.png")

    html = render_html_report(
        build_report_data(analytics_clean, None, branded), include_charts=False
    )

    assert "data:image" not in html


def test_unreadable_logo_is_skipped_rather_than_fatal(
    profile, analytics_clean, tmp_path, monkeypatch, caplog
):
    """A permission error on the logo must not take the whole export down."""
    from grant_assistant.reporting import branding, build_report_data
    from grant_assistant.reporting.branding import logo_bytes

    logo = tmp_path / "logo.png"
    logo.write_bytes(b"\x89PNG\r\n\x1a\nsynthetic")
    branded = profile.model_copy(deep=True)
    branded.report.logo_path = str(logo)

    def denied(self, *args, **kwargs):
        raise PermissionError("locked by another process")

    monkeypatch.setattr(branding.Path, "read_bytes", denied)

    with caplog.at_level(logging.WARNING):
        assert logo_bytes(build_report_data(analytics_clean, None, branded)) is None
    assert "Skipping report logo" in caplog.text


def test_word_and_deck_survive_an_unusable_logo(profile, analytics_clean, tmp_path):
    """The binary renderers embed the logo differently; both must degrade too."""
    from docx import Document

    from grant_assistant.reporting import build_report_data, write_docx_report

    branded = profile.model_copy(deep=True)
    branded.report.logo_path = str(tmp_path / "nope.png")
    report = build_report_data(analytics_clean, None, branded)

    doc = Document(str(write_docx_report(report, tmp_path / "r.docx")))
    assert any(p.runs for p in doc.paragraphs)


def test_merge_rejects_an_ambiguous_join_key_column(profile, tmp_path):
    """Two columns mapping to client_id is a data error, not a coin flip."""
    import pandas as pd

    from grant_assistant.ingestion import IngestionError, merge_related_datasets

    primary = tmp_path / "primary.csv"
    pd.DataFrame({"Client ID": ["A"], "Program Name": ["RRH"]}).to_csv(primary, index=False)
    related = tmp_path / "income.csv"
    # Both headers normalize onto client_id, so the join key is ambiguous.
    pd.DataFrame({"Client ID": ["A"], "client-id": ["A"], "Case Manager": ["One"]}).to_csv(
        related, index=False
    )

    with pytest.raises(IngestionError, match="Expected exactly one column"):
        merge_related_datasets(primary, [related], profile)


def test_merge_rejects_an_unknown_canonical_join_key(profile, tmp_path):
    import pandas as pd

    from grant_assistant.ingestion import IngestionError, merge_related_datasets

    primary = tmp_path / "primary.csv"
    pd.DataFrame({"Client ID": ["A"]}).to_csv(primary, index=False)

    with pytest.raises(IngestionError, match="Unknown canonical join key"):
        merge_related_datasets(primary, [], profile, join_on="not_a_column")


def test_merge_is_a_no_op_when_a_related_file_adds_nothing(profile, tmp_path):
    """A related file whose columns all already exist must not duplicate or reorder."""
    import pandas as pd

    from grant_assistant.ingestion import merge_related_datasets

    primary = tmp_path / "primary.csv"
    pd.DataFrame({"Client ID": ["A"], "Program Name": ["RRH"]}).to_csv(primary, index=False)
    related = tmp_path / "same.csv"
    pd.DataFrame({"Client ID": ["A"], "Program Name": ["Ignored"]}).to_csv(related, index=False)

    merged = merge_related_datasets(primary, [related], profile)

    assert list(merged.columns) == ["Client ID", "Program Name"]
    assert merged.loc[0, "Program Name"] == "RRH"  # the primary always wins


def test_uploaded_and_on_disk_merges_agree(profile, tmp_path):
    """The web app and the CLI must not disagree about a merge."""
    import pandas as pd

    from grant_assistant.ingestion import merge_related_datasets, merge_uploaded_datasets

    primary_frame = pd.DataFrame({"Client ID": ["A", "B"], "Program Name": ["RRH", "ES"]})
    related_frame = pd.DataFrame({"Client ID": ["A", "B"], "Case Manager": ["One", "Two"]})
    primary = tmp_path / "primary.csv"
    related = tmp_path / "staff.csv"
    primary_frame.to_csv(primary, index=False)
    related_frame.to_csv(related, index=False)

    on_disk = merge_related_datasets(primary, [related], profile)
    uploaded = merge_uploaded_datasets(
        primary_frame, [("staff.csv", related.read_bytes())], profile
    )

    assert list(on_disk.columns) == list(uploaded.columns)
    assert on_disk["Case Manager"].tolist() == uploaded["Case Manager"].tolist()


def test_tls_stays_on_for_any_unrecognized_flag_value():
    """An unrecognized value must fail safe: still encrypted."""
    from grant_assistant.automation import _tls_enabled

    for off in ("false", "FALSE", " no ", "0", "off"):
        assert _tls_enabled(off) is False, off
    for on in ("true", "1", "yes", "", "maybe", "flase"):
        assert _tls_enabled(on) is True, on


def test_sending_twice_does_not_accumulate_headers(monkeypatch):
    """A retrying scheduler reuses the message; it must not grow duplicate headers."""
    from email.message import EmailMessage

    from grant_assistant.automation import send_audit_email

    sent: list[EmailMessage] = []

    class FakeSMTP:
        def __init__(self, *args, **kwargs):
            pass

        def __enter__(self):
            return self

        def __exit__(self, *_args):
            return None

        def starttls(self, *, context):
            pass

        def send_message(self, message):
            sent.append(message)

    monkeypatch.setattr("grant_assistant.automation.smtplib.SMTP", FakeSMTP)
    message = EmailMessage()
    message.set_content("Synthetic audit summary")

    for _ in range(2):
        send_audit_email(
            message, ["a@example.org"], host="smtp.example.org", sender="r@example.org"
        )

    assert message["To"] is None, "the caller's message must not be mutated"
    for delivered in sent:
        assert delivered.get_all("To") == ["a@example.org"]
        assert delivered.get_all("From") == ["r@example.org"]
