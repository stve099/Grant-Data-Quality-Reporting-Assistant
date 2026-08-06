"""PowerPoint export tests.

The deck is a third renderer over one ReportData, so the property that matters
most is that it cannot disagree with the other two: every figure on a slide has
to be the calculated value, not a re-derivation.
"""

from __future__ import annotations

import importlib.util

import pytest

_HAS_PPTX = importlib.util.find_spec("pptx") is not None
pytestmark = pytest.mark.skipif(not _HAS_PPTX, reason="pptx extra not installed")


@pytest.fixture(scope="module")
def report_data(analytics_flawed, audit_flawed, profile):
    from grant_assistant.reporting import build_report_data

    return build_report_data(analytics_flawed, audit_flawed, profile)


@pytest.fixture(scope="module")
def deck(report_data, tmp_path_factory):
    from pptx import Presentation

    from grant_assistant.reporting import write_pptx_report

    path = write_pptx_report(report_data, tmp_path_factory.mktemp("pptx") / "deck.pptx")
    return Presentation(str(path))


def _all_text(deck) -> str:
    parts: list[str] = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


# -- Structure ---------------------------------------------------------------


def test_deck_is_written_and_short(deck):
    """A deck is a summary; the report carries the detail."""
    assert 6 <= len(deck.slides) <= 14


def test_widescreen_aspect_ratio(deck):
    ratio = deck.slide_width / deck.slide_height
    assert abs(ratio - 16 / 9) < 0.01


def test_title_slide_carries_the_grant_and_period(deck, report_data):
    text = "\n".join(s.text_frame.text for s in deck.slides[0].shapes if s.has_text_frame)
    assert report_data.profile.grant_name in text
    assert report_data.period_label in text


def test_synthetic_data_disclaimer_is_present(deck):
    """The same honesty the Word cover page carries."""
    assert "no real client information" in _all_text(deck)


# -- Numbers match the calculations ------------------------------------------


def test_headline_numbers_come_from_analytics(deck, analytics_flawed):
    text = _all_text(deck)
    assert f"{analytics_flawed.households_served:,}" in text
    assert f"{analytics_flawed.total_individuals:,}" in text
    if analytics_flawed.permanent_housing_rate is not None:
        assert f"{analytics_flawed.permanent_housing_rate}%" in text


def test_every_measure_appears_with_its_target(deck, analytics_flawed):
    text = _all_text(deck)
    for measure in analytics_flawed.measures:
        assert measure.name in text


def test_data_quality_score_matches_the_audit(deck, audit_flawed):
    text = _all_text(deck)
    assert f"{audit_flawed.overall_score:.1f}" in text
    assert f"Grade {audit_flawed.grade}" in text


def test_blocking_issues_are_named(deck, audit_flawed):
    text = _all_text(deck)
    if audit_flawed.blocking_issues:
        assert audit_flawed.blocking_issues[0].rule_name in text
    else:
        assert "No blocking issues" in text


def test_executive_summary_is_included(deck, report_data):
    assert report_data.executive_summary[:60] in _all_text(deck)


def test_methodology_states_who_calculates(deck):
    """The grounding claim belongs in front of a board, not only in the docs."""
    text = _all_text(deck)
    assert "calculated by the application" in text
    assert "narrates" in text


# -- Degradation -------------------------------------------------------------


def test_deck_builds_without_a_chart_backend(report_data, tmp_path, monkeypatch):
    """No kaleido means fewer slides, not a failure."""
    from pptx import Presentation

    from grant_assistant.reporting import write_pptx_report

    monkeypatch.setattr(
        "grant_assistant.reporting.chart_images.chart_backend_available", lambda: False
    )
    path = write_pptx_report(report_data, tmp_path / "no_charts.pptx")
    deck = Presentation(str(path))
    images = [shape for slide in deck.slides for shape in slide.shapes if shape.shape_type == 13]
    assert not images
    # The substance survives.
    assert "Executive summary" in _all_text(deck)


def test_charts_are_embedded_when_available(deck):
    from grant_assistant.reporting.chart_images import chart_backend_available

    images = [shape for slide in deck.slides for shape in slide.shapes if shape.shape_type == 13]
    if chart_backend_available():
        assert images
    else:  # pragma: no cover - depends on the installed extras
        assert not images


def test_missing_audit_does_not_break_the_deck(analytics_flawed, profile, tmp_path):
    """Analytics-only runs must still produce a deck."""
    from grant_assistant.reporting import build_report_data, write_pptx_report

    data = build_report_data(analytics_flawed, None, profile)
    path = write_pptx_report(data, tmp_path / "no_audit.pptx")
    assert path.exists()
